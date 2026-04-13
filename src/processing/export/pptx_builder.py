from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.util import Pt

from src.memory.wip.database import WIPDatabase
from src.memory.wip.schema import BulletPoint, ProtoSlide, SlideContent

# ── Layout name constants ─────────────────────────────────────────────────────
# These match the built-in layout names in the default python-pptx blank
# template.  Index fallbacks are used when a name match fails so the builder
# never hard-crashes on a custom template with renamed layouts.
_LAYOUT_TITLE_SLIDE = "Title Slide"
_LAYOUT_TITLE_AND_CONTENT = "Title and Content"
_LAYOUT_TWO_CONTENT = "Two Content"

# Font overrides
_BIG_NUMBER_FONT_PT = 40
_QUOTE_FONT_PT = 20


def _split_bold_runs(text: str, bold_phrases: List[str]) -> List[tuple[str, bool]]:
    """
    Splits *text* into a list of (segment, is_bold) tuples.

    Each phrase in *bold_phrases* that appears as a substring of *text* is
    extracted as a bold segment.  Matching is left-to-right; overlapping
    intervals are resolved by taking the earlier one.  Phrases absent from
    *text* are silently skipped so a stale or hallucinated phrase never
    crashes the builder.
    """
    if not bold_phrases:
        return [(text, False)]

    intervals: List[tuple[int, int]] = []
    for phrase in bold_phrases:
        if not phrase:
            continue
        idx = text.find(phrase)
        if idx == -1:
            continue
        intervals.append((idx, idx + len(phrase)))

    if not intervals:
        return [(text, False)]

    intervals.sort()
    merged: List[tuple[int, int]] = []
    for start, end in intervals:
        if merged and start < merged[-1][1]:
            continue  # overlaps the previous interval — skip
        merged.append((start, end))

    runs: List[tuple[str, bool]] = []
    cursor = 0
    for start, end in merged:
        if cursor < start:
            runs.append((text[cursor:start], False))
        runs.append((text[start:end], True))
        cursor = end
    if cursor < len(text):
        runs.append((text[cursor:], False))

    return runs


class PptxBuilder:
    """
    Reads all ProtoSlide rows from a WIPDatabase and renders them into a
    .pptx file at the given output path.

    Usage::

        with WIPDatabase() as db:
            PptxBuilder(output_path=Path("output.pptx"), db=db).build()
    """

    # Maps SlideContent.layout literal → (pptx layout name, special mode)
    _LAYOUT_MAP = {
        "title_slide":    (_LAYOUT_TITLE_SLIDE,       None),
        "title_and_body": (_LAYOUT_TITLE_AND_CONTENT, None),
        "two_column":     (_LAYOUT_TWO_CONTENT,       None),
        "big_number":     (_LAYOUT_TITLE_AND_CONTENT, "big_number"),
        "quote":          (_LAYOUT_TITLE_AND_CONTENT, "quote"),
        "media_left":     (_LAYOUT_TWO_CONTENT,       "media_placeholder"),
        "media_right":    (_LAYOUT_TWO_CONTENT,       "media_placeholder"),
    }

    def __init__(self, output_path: Path, db: WIPDatabase) -> None:
        self.output_path = Path(output_path)
        self._db = db

    # ── Public API ────────────────────────────────────────────────────────────

    def build(self) -> Path:
        """
        Loads all proto-slides from wip.db in slide-number order and writes
        the presentation to self.output_path.  Returns the resolved path.

        Raises:
            ValueError: if no slides exist in the database.
        """
        slide_numbers = self._db.list_slide_numbers()
        if not slide_numbers:
            raise ValueError(
                "PptxBuilder: no proto-slides found in wip.db. "
                "Run the slide-generation graph first."
            )

        slides: List[ProtoSlide] = [
            self._db.load_slide(n) for n in slide_numbers
        ]

        prs = Presentation()
        for slide in slides:
            if slide is not None:
                self._add_slide(prs, slide)

        self.output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(self.output_path))
        return self.output_path.resolve()

    # ── Private helpers ───────────────────────────────────────────────────────

    def _add_slide(self, prs: Presentation, proto: ProtoSlide) -> None:
        content: SlideContent = proto.content
        layout_name, mode = self._LAYOUT_MAP.get(
            content.layout, (_LAYOUT_TITLE_AND_CONTENT, None)
        )
        layout = self._pick_layout(prs, layout_name)
        slide = prs.slides.add_slide(layout)

        self._set_title(slide, content.title)
        self._set_body(slide, content, mode)
        self._set_speaker_notes(slide, content.speaker_notes)

    def _pick_layout(self, prs: Presentation, name: str):
        """Returns the layout matching *name*, falling back to index 1."""
        for layout in prs.slide_layouts:
            if layout.name == name:
                return layout
        # Graceful fallback: index 0 = Title Slide, 1 = Title and Content
        fallback_idx = 0 if name == _LAYOUT_TITLE_SLIDE else 1
        return prs.slide_layouts[fallback_idx]

    def _set_title(self, slide, title: str) -> None:
        title_shape = slide.shapes.title
        if title_shape is not None:
            title_shape.text = title

    def _set_body(self, slide, content: SlideContent, mode: str | None) -> None:
        """Finds the body/content placeholder and writes bullets into it."""
        body_ph = self._find_body_placeholder(slide)
        if body_ph is None:
            return

        tf = body_ph.text_frame
        tf.word_wrap = True

        first = True
        for bullet in content.bullets:
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            para.level = 0
            self._write_run(para, bullet.text, bullet, mode)

            for sub in bullet.sub_bullets:
                sub_para = tf.add_paragraph()
                sub_para.level = 1
                run = sub_para.add_run()
                run.text = sub
                if mode == "quote":
                    run.font.italic = True

    def _find_body_placeholder(self, slide):
        """
        Returns the first non-title placeholder on the slide, which is the
        content/body area in all standard layouts.
        """
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:
                return ph
        return None

    def _write_run(self, para, text: str, bullet: BulletPoint, mode: str | None) -> None:
        """
        Splits *text* into plain/bold segments (via bold_phrases) and writes
        one python-pptx run per segment.  Mode-level styling (big_number, quote)
        is applied to every run; bold_phrases styling is applied only to the
        matched segments.
        """
        segments = _split_bold_runs(text, bullet.bold_phrases)

        for segment_text, is_bold in segments:
            run = para.add_run()
            run.text = segment_text

            if mode == "big_number":
                run.font.size = Pt(_BIG_NUMBER_FONT_PT)
                run.font.bold = True
            elif mode == "quote":
                run.font.italic = True
                run.font.size = Pt(_QUOTE_FONT_PT)

            if is_bold:
                run.font.bold = True

    def _set_speaker_notes(self, slide, notes: str) -> None:
        if not notes:
            return
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
